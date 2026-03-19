"""
ND-CORE Flask 백엔드 서버
- /search   : 파일 이름 검색
- /open_explorer : 파일 탐색기에서 경로 열기
"""

import os
import subprocess
import math
from datetime import datetime
from flask import Flask, request, jsonify, send_from_directory
from flask_cors import CORS

app = Flask(__name__, static_folder='frontend', static_url_path='/frontend')
CORS(app)  # 모든 origin 허용 (로컬 HTML 파일에서의 fetch 허용)
os.environ["PADDLE_PDX_DISABLE_MODEL_SOURCE_CHECK"] = "True"
try:
    from services.ai_service import AIService
    ai_service = AIService()
except Exception as e:
    print("AI 서비스(Chroma 및 Gemini 임베딩) 초기화 실패:", e)
    ai_service = None

try:
    import docx
except ImportError:
    docx = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

try:
    from paddleocr import PaddleOCR
    ocr = PaddleOCR(use_textline_orientation=True, lang='korean')
except ImportError:
    ocr = None

# ─────────────────────────────────────────
#  검색할 루트 경로 설정 (필요 시 변경)
# ─────────────────────────────────────────
SEARCH_ROOT = "C:\\"   # 전체 C 드라이브 검색. 특정 폴더로 좁히려면 변경하세요.
MAX_RESULTS = 200       # 최대 반환 개수


def format_size(size_bytes: int) -> str:
    """파일 크기를 읽기 쉬운 단위로 변환"""
    if size_bytes == 0:
        return "0 B"
    units = ["B", "KB", "MB", "GB", "TB"]
    i = int(math.floor(math.log(size_bytes, 1024)))
    p = math.pow(1024, i)
    s = round(size_bytes / p, 2)
    return f"{s} {units[i]}"


def format_modified(timestamp: float) -> str:
    """타임스탬프를 날짜 문자열로 변환"""
    try:
        dt = datetime.fromtimestamp(timestamp)
        return dt.strftime("%Y-%m-%d %H:%M")
    except Exception:
        return "-"


def search_files_content(query: str, search_root: str, max_results: int):
    """
    파일의 내용을 기반으로 텍스트 검색을 수행하는 임시 시멘틱 검색 함수입니다.
    """
    results = []
    query_lower = query.lower()

    SKIP_DIRS = {
        "$recycle.bin", "system volume information", "windows",
        "windowsapps", "program files (x86)", "appdata",
        "$windows.~bt", "$windows.~ws",
    }

    # 텍스트 검색을 시도할 확장자 목록
    TEXT_EXTS = {".txt", ".md", ".csv", ".json", ".py", ".js", ".html", ".css", ".xml"}
    IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".bmp"}
    OFFICE_EXTS = {".docx", ".pptx", ".xlsx", ".pdf"}

    for dirpath, dirnames, filenames in os.walk(search_root, topdown=True, onerror=lambda e: None):
        dirnames[:] = [d for d in dirnames if d.lower() not in SKIP_DIRS and not d.startswith('.')]
        for filename in filenames:
            ext = os.path.splitext(filename)[1].lower()
            if ext in TEXT_EXTS or ext in OFFICE_EXTS or ext in IMAGE_EXTS:
                full_path = os.path.join(dirpath, filename)
                try:
                    lines = []
                    # 파일 형식에 따른 텍스트 추출 로직
                    if ext in TEXT_EXTS:
                        with open(full_path, 'r', encoding='utf-8', errors='ignore') as f:
                            lines = f.readlines()
                    elif ext == '.docx' and docx is not None:
                        doc = docx.Document(full_path)
                        lines = [p.text for p in doc.paragraphs if p.text.strip()]
                    elif ext == '.pptx' and Presentation is not None:
                        prs = Presentation(full_path)
                        for slide in prs.slides:
                            for shape in slide.shapes:
                                if hasattr(shape, "text") and shape.text.strip():
                                    lines.extend(shape.text.split('\n'))
                    elif ext == '.xlsx' and openpyxl is not None:
                        wb = openpyxl.load_workbook(full_path, data_only=True)
                        for sheet in wb.worksheets:
                            for row in sheet.iter_rows(values_only=True):
                                row_texts = [str(cell) for cell in row if cell is not None]
                                if row_texts:
                                    lines.append(" ".join(row_texts))
                        wb.close()
                    elif ext == '.pdf' and PyPDF2 is not None:
                        with open(full_path, 'rb') as f:
                            reader = PyPDF2.PdfReader(f)
                            for page in reader.pages:
                                text = page.extract_text()
                                if text:
                                    lines.extend(text.split('\n'))
                    elif ext in IMAGE_EXTS and ocr is not None:
                        result = ocr.ocr(full_path, cls=True)
                        if result and result[0]:
                            for line_data in result[0]:
                                lines.append(line_data[1][0])

                    if not lines:
                        continue

                    match_idx = -1
                    for idx, line in enumerate(lines):
                        if query_lower in line.lower():
                            match_idx = idx
                            break

                    if match_idx != -1:
                        start_idx = max(0, match_idx - 1)
                        end_idx = min(len(lines), match_idx + 2)
                        snippet_lines = [l.strip() for l in lines[start_idx:end_idx] if l.strip()]
                        snippet = "\n".join(snippet_lines)

                        stat = os.stat(full_path)
                        results.append({
                            "name": filename,
                            "path": full_path,
                            "size": format_size(stat.st_size),
                            "modified": format_modified(stat.st_mtime),
                            "snippet": snippet
                        })
                        if len(results) >= max_results:
                            return results
                except Exception:
                    pass

    return results

def search_files(query: str, search_root: str, max_results: int, search_type: str = "name"):
    """
    지정된 루트 경로에서 query를 포함하는 파일을 재귀적으로 탐색합니다.
    """
    results = []
    query_lower = query.lower()

    # 접근 불가 폴더 스킵 목록
    SKIP_DIRS = {
        "$recycle.bin", "system volume information", "windows",
        "windowsapps", "program files (x86)", "appdata",
        "$windows.~bt", "$windows.~ws",
    }

    for dirpath, dirnames, filenames in os.walk(search_root, topdown=True, onerror=lambda e: None):
        # 스킵할 폴더 제거
        dirnames[:] = [
            d for d in dirnames
            if d.lower() not in SKIP_DIRS and not d.startswith('.')
        ]

        for filename in filenames:
            match = False
            if search_type == "name":
                ext = os.path.splitext(filename)[1].lower().strip(".")
                query_ext = query_lower.strip(".")
                # 파일명 포함되거나, 확장자가 정확히 일치하거나
                if ext == query_ext or query_lower in filename.lower():
                    match = True
            else:
                if query_lower in filename.lower():
                    match = True
            
            if match:
                full_path = os.path.join(dirpath, filename)
                try:
                    stat = os.stat(full_path)
                    results.append({
                        "name": filename,
                        "path": full_path,
                        "size": format_size(stat.st_size),
                        "modified": format_modified(stat.st_mtime),
                    })
                except (PermissionError, FileNotFoundError):
                    pass

                if len(results) >= max_results:
                    return results

    return results


@app.route("/")
def index():
    """메인 UI 페이지 제공"""
    return send_from_directory(app.root_path, 'search_ui.html')


@app.route("/search", methods=["POST"])
def search():
    """
    파일 검색 API
    Request body: { "query": "검색어", "searchType": "name" }
    Response: { "slots": [{ "file_info": { name, path, size, modified } }] }
    """
    data = request.get_json(force=True)
    query = (data.get("query") or "").strip()
    search_type = data.get("searchType", "name")

    if not query:
        return jsonify({"error": "검색어를 입력하세요.", "slots": []}), 400

    if search_type == "content":
        try:
            if ai_service is not None:
                # 임베딩 & Chroma 활용하는 RAG 검색 수행
                files = ai_service.search(query, MAX_RESULTS)
            else:
                # fallback (AI 모듈 에러 시 텍스트 검색 사용)
                files = search_files_content(query, SEARCH_ROOT, MAX_RESULTS) 
        except Exception as e:
            print(f"시멘틱 검색 오류: {e}")
            return jsonify({"error": f"시멘틱 검색 오류: {e}", "slots": []}), 500
    else:
        files = search_files(query, SEARCH_ROOT, MAX_RESULTS, search_type)

    slots = [{"file_info": f} for f in files]
    return jsonify({
        "query": query,
        "total": len(slots),
        "slots": slots
    })


@app.route("/index_folder", methods=["POST"])
def index_folder():
    """
    RAG용 벡터 DB 생성을 위해 특정 폴더를 스캔하여 텍스트를 파싱하고 Chroma DB에 저장합니다.
    Request body: { "folder_path": "C:\\..." }
    """
    if ai_service is None:
        return jsonify({"error": "AI 서비스 모듈이 로드되지 않았습니다. (패키지 또는 API 키 확인 필요)"}), 500

    data = request.get_json(force=True)
    folder = (data.get("folder_path") or "").strip()
    
    if not folder or not os.path.exists(folder):
        return jsonify({"error": "유효하지 않은 경로입니다."}), 400

    try:
        res = ai_service.index_directory(folder)
        if "error" in res:
            return jsonify(res), 400
        return jsonify(res)
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/open_explorer", methods=["POST"])
def open_explorer():
    """
    파일 탐색기에서 파일 경로를 열기
    Request body: { "path": "C:\\..." }
    """
    data = request.get_json(force=True)
    path = (data.get("path") or "").strip()

    if not path or not os.path.exists(path):
        return jsonify({"error": "유효하지 않은 경로입니다."}), 400

    try:
        # 파일이면 해당 폴더를 열면서 파일 선택
        if os.path.isfile(path):
            subprocess.Popen(f'explorer /select,"{path}"')
        else:
            subprocess.Popen(f'explorer "{path}"')
        return jsonify({"status": "ok"})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/health", methods=["GET"])
def health():
    return jsonify({"status": "running", "search_root": SEARCH_ROOT})


if __name__ == "__main__":
    print("=" * 50)
    print("  ND-CORE Flask 서버 시작")
    print(f"  검색 루트: {SEARCH_ROOT}")
    print(f"  주소: http://127.0.0.1:5000")
    print("=" * 50)
    
    # 서버 기동 시 지정된 경로의 파일들을 자동으로 벡터 DB에 인덱싱 (RAG용)
    if ai_service is not None:
        target_folder = r"C:\PJT_new\educate"
        print(f"[*] 서버 시작 중... RAG 파이프라인 자동 인덱싱 진행: {target_folder}")
        try:
            res = ai_service.index_directory(target_folder)
            print("#####################")
            print(f"res = {res}")
            print("#####################")
            print("[*] 인덱싱 완료:", res)
        except Exception as e:
            print(f"[*] 자동 인덱싱 실패: {e}")
            
    app.run(host="127.0.0.1", port=5000, debug=False)

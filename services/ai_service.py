import sys
print("#####################")
print(sys.executable)
print("#####################")

import os
import math
from datetime import datetime
from uuid import uuid4
from langchain_text_splitters import RecursiveCharacterTextSplitter
# 설정 예시: 텍스트를 1000자 단위로 자르고, 문맥 유지를 위해 100자씩 겹치게 설정
text_splitter = RecursiveCharacterTextSplitter(
    chunk_size=1000,
    chunk_overlap=100
)

try:
    from langchain_google_genai import GoogleGenerativeAIEmbeddings
    from langchain_chroma import Chroma
    from langchain_google_genai import GoogleGenerativeAIEmbeddings
    from langchain_chroma import Chroma
    from langchain_core.documents import Document

    # 초기화는 __init__에서 진행하므로 전역 변수 초기화 코드는 삭제했습니다.

    print("#####################")
    print(f"Chroma = {Chroma}")
    print(f"GoogleGenerativeAIEmbeddings = {GoogleGenerativeAIEmbeddings}")
    print("#####################")

except ImportError:
    GoogleGenerativeAIEmbeddings = None
    Chroma = None

try:
    import docx
except ImportError:
    docx = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

try:
    import sys
    class _DummyLangchain: pass
    sys.modules.setdefault('langchain.docstore', _DummyLangchain())
    import langchain_core.documents
    sys.modules.setdefault('langchain.docstore.document', langchain_core.documents)
    import langchain_text_splitters
    sys.modules.setdefault('langchain.text_splitter', langchain_text_splitters)
    from paddleocr import PaddleOCR
    ocr = PaddleOCR(use_textline_orientation=True, lang='korean')
except ImportError:
    ocr = None


class AIService:
    def __init__(self, db_path=r"C:\PJT_new\DB"):
        from dotenv import load_dotenv
        load_dotenv(r"C:\PJT_new\backend\.env")

        if Chroma is None or GoogleGenerativeAIEmbeddings is None:
            raise ImportError("langchain-google-genai, langchain-chroma 패키지가 필요합니다.")

        # Gemini API 키 확인
        api_key = "AIzaSyAjtzUkgtStQToyx_3rpkluGIhn84xabkQ"
        # api_key = os.environ.get("GOOGLE_API_KEY") or os.environ.get("GEMINI_API_KEY")
        if not api_key:
            print("경고: GOOGLE_API_KEY 또는 GEMINI_API_KEY 환경변수가 설정되지 않아 임베딩 모델 로드에 실패할 수 있습니다.")
        else:
            os.environ["GOOGLE_API_KEY"] = api_key

        print("[AIService] Gemini Embedding 및 Chroma Vector DB 초기화 중...")
        
        # 임베딩 모델: Gemini Text Embedding 004
        self.embeddings = GoogleGenerativeAIEmbeddings(model="models/text-embedding-004")
        self.db_path = db_path
        
        # 벡터 디비 (Chroma)
        self.vector_store = Chroma(
            collection_name="nd_core_files",
            embedding_function=self.embeddings,
            persist_directory=self.db_path
        )
        
        # 청크 분할기 (텍스트를 의미 단위로 자름)
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            separators=["\n\n", "\n", ".", " ", ""]
        )
        print(f"[AIService] Chroma DB 준비 완료: {self.db_path}")

    def _format_size(self, size_bytes: int) -> str:
        if size_bytes == 0:
            return "0 B"
        units = ["B", "KB", "MB", "GB", "TB"]
        i = int(math.floor(math.log(size_bytes, 1024)))
        p = math.pow(1024, i)
        s = round(size_bytes / p, 2)
        return f"{s} {units[i]}"

    def _format_modified(self, timestamp: float) -> str:
        try:
            dt = datetime.fromtimestamp(timestamp)
            return dt.strftime("%Y-%m-%d %H:%M")
        except Exception:
            return "-"

    def search(self, query: str, max_results: int = 5):
        """저장된 벡터 DB(Chroma)에서 질문과 텍스트 의미가 유사한 문서를 반환 (RAG의 Retrieval 단계)"""
        print(f"[AIService] '{query}' 벡터 유사도 검색 시작...")
        
        # k=max_results 개수만큼 관련성이 높은 문서 청크 반환
        docs = self.vector_store.similarity_search_with_relevance_scores(query, k=max_results)
        
        results = []
        for doc, score in docs:
            full_path = doc.metadata.get("source", "Unknown")
            filename = os.path.basename(full_path)
            
            try:
                stat = os.stat(full_path)
                size_str = self._format_size(stat.st_size)
                mod_str = self._format_modified(stat.st_mtime)
            except Exception:
                size_str = "-"
                mod_str = "-"
                
            # 검색된 텍스트 청크를 스니펫으로 반환
            snippet = doc.page_content.strip()
            
            results.append({
                "name": filename,
                "path": full_path,
                "size": size_str,
                "modified": mod_str,
                "snippet": snippet,
                "similarity": score
            })
            
        print(f"[AIService] 검색 완료! {len(results)}건 반환됨.")
        return results

    def index_directory(self, folder_path: str):
        """특정 폴더의 파일들을 순회하며 내용을 추출하고 벡터화하여 DB에 삽입 (RAG 데이터 파이프라인)"""
        if not os.path.exists(folder_path):
            return {"error": "경로가 존재하지 않습니다."}

        TEXT_EXTS = {".txt", ".md", ".csv", ".json", ".py", ".js", ".html", ".css", ".xml"}
        IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".bmp"}
        OFFICE_EXTS = {".docx", ".pptx", ".xlsx", ".pdf"}
        
        SKIP_DIRS = {
            "$recycle.bin", "system volume information", "windows",
            "windowsapps", "program files (x86)", "appdata",
            "$windows.~bt", "$windows.~ws",
        }

        docs_to_index = []
        
        print(f"[AIService] {folder_path} 경로의 파일 인덱싱 시작...")
        
        for dirpath, dirnames, filenames in os.walk(folder_path):
            dirnames[:] = [d for d in dirnames if d.lower() not in SKIP_DIRS and not d.startswith('.')]
            
            for filename in filenames:
                ext = os.path.splitext(filename)[1].lower()
                full_path = os.path.join(dirpath, filename)
                
                if ext in TEXT_EXTS or ext in OFFICE_EXTS or ext in IMAGE_EXTS:
                    text_content = ""
                    
                    try:
                        # 텍스트 추출 로직 (app.py와 동일하게 처리)
                        if ext in TEXT_EXTS:
                            with open(full_path, 'r', encoding='utf-8', errors='ignore') as f:
                                text_content = f.read()
                        elif ext == '.docx' and docx is not None:
                            doc = docx.Document(full_path)
                            text_content = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
                        elif ext == '.pptx' and Presentation is not None:
                            prs = Presentation(full_path)
                            for slide in prs.slides:
                                for shape in slide.shapes:
                                    if hasattr(shape, "text") and shape.text.strip():
                                        text_content += shape.text + "\n"
                        elif ext == '.xlsx' and openpyxl is not None:
                            wb = openpyxl.load_workbook(full_path, data_only=True)
                            for sheet in wb.worksheets:
                                for row in sheet.iter_rows(values_only=True):
                                    row_texts = [str(cell) for cell in row if cell is not None]
                                    if row_texts:
                                        text_content += " ".join(row_texts) + "\n"
                            wb.close()
                        elif ext == '.pdf' and PyPDF2 is not None:
                            with open(full_path, 'rb') as f:
                                reader = PyPDF2.PdfReader(f)
                                for page in reader.pages:
                                    p_text = page.extract_text()
                                    if p_text:
                                        text_content += p_text + "\n"
                        elif ext in IMAGE_EXTS and ocr is not None:
                            result = ocr.ocr(full_path, cls=True)
                            if result and result[0]:
                                for line_data in result[0]:
                                    text_content += line_data[1][0] + "\n"
                                    
                    except Exception as e:
                        print(f"[{filename}] 추출 중 오류: {e}")
                        
                    if text_content.strip():
                        print(f"추출시작")
                        # 추출된 텍스트를 청크 분할 (Chunking)
                        chunks = self.text_splitter.split_text(text_content)
                        for chunk in chunks:
                            docs_to_index.append(
                                Document(
                                    page_content=chunk,
                                    metadata={"source": full_path, "filename": filename}
                                )
                            )
        print("###############")
        print(f"docs_to_index = {docs_to_index}")
        print("###############")
        if docs_to_index:
            print(f"[AIService] 총 {len(docs_to_index)}개의 청크가 생성되었습니다. 벡터 DB 저장 중...")
            # ChromaDB에 청크와 메타데이터 삽입 (임베딩 변환 및 저장됨)
            self.vector_store.add_documents(documents=docs_to_index)
            print("[AIService] 인덱싱(벡터 변환 및 DB 저장)이 성공적으로 완료되었습니다!")
            return {"status": "success", "chunks_added": len(docs_to_index)}
        else:
            print("[AIService] 추출된 텍스트가 없어 저장할 내용이 없습니다.")
            return {"status": "empty", "chunks_added": 0}
